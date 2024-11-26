from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import os

def extract_text(file_path):
    """
    Extracts text from a file based on its format.
    :param file_path: Path to the file.
    :return: Extracted text as a string.
    """
    file_path = file_path.lower()

    if file_path.endswith('.pdf'):
        return extract_pdf_text(file_path)
    elif file_path.endswith('.docx'):
        return extract_word_text(file_path)
    elif file_path.endswith('.pptx'):
        return extract_ppt_text(file_path)
    elif file_path.endswith(('.png', '.jpg', '.jpeg')):
        return extract_image_text(file_path)
    else:
        raise ValueError(f"Unsupported file format for file: {file_path}")

def extract_pdf_text(file_path):
    """
    Extracts text from a PDF file.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        reader = PdfReader(file_path)
        return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    except Exception as e:
        raise ValueError(f"Error processing PDF file {file_path}: {str(e)}")

def extract_word_text(file_path):
    """
    Extracts text from a Word document.
    """
    doc = Document(file_path)
    return " ".join([p.text for p in doc.paragraphs])

def extract_ppt_text(file_path):
    """
    Extracts text from a PowerPoint presentation.
    """
    pres = Presentation(file_path)
    return " ".join([shape.text for slide in pres.slides for shape in slide.shapes if shape.has_text_frame])

def extract_image_text(file_path):
    """
    Extracts text from an image file.
    """
    try:
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
    except Exception as e:
        raise ValueError(f"Error processing image file {file_path}: {str(e)}")
